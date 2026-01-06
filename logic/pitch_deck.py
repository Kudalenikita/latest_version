# logic/pitch_deck.py
import json
import asyncio
import os
import re
from datetime import datetime

from autogen_agentchat.agents import AssistantAgent
from rag.rag_engine import query_vector_db
from db.db_utils import load_contracts_for_customer, load_all_releases_for_customer


async def generate_pitch_deck_content(
    agent: AssistantAgent,
    customer_name: str,
    vector_client,
    embedding_func,
    comparison,
    risk_data
) -> dict:
    """
    Generate structured pitch deck content using structured JSON output.
    Returns: dict with slide1_title, slide1_content, ..., slide7_title, slide7_content
    """
    contracts_df = load_contracts_for_customer(customer_name)
    releases_df = load_all_releases_for_customer(customer_name)

    total = len(contracts_df["feature_id"].unique()) if not contracts_df.empty else 0
    released = planned = missing = 0

    if not releases_df.empty:
        status_lower = releases_df["status"].astype(str).str.lower()
        released = len(releases_df[status_lower.isin(["released", "done", "completed", "live"])]["feature_id"].unique())
        planned = len(releases_df[status_lower == "planned"]["feature_id"].unique())
        missing = total - released - planned

    data_summary = f"""Customer: {customer_name}
Total committed features: {total}
Delivered: {released}
Planned / In progress: {planned}
Not yet addressed: {missing}"""

    # Get relevant context from vector DB
    rag_docs = query_vector_db(
        vector_client,
        embedding_func,
        "strategic overview contract commitments roadmap delivery status risks value proposition",
        customer_filter=customer_name,
        n_results=12
    )
    rag_context = "\n\n".join(doc["text"] for doc in rag_docs) if rag_docs else ""

    task = f"""You are a professional sales strategist generating a 7-slide pitch deck in strict JSON format.

Respond with ONLY a valid JSON object. No explanations, no markdown, no code fences, no extra text.

Use exactly these keys:
{{
  "slide1_title": "string",
  "slide1_content": "string (use \\n for line breaks)",
  "slide2_title": "string",
  "slide2_content": "string",
  "slide3_title": "string",
  "slide3_content": "string",
  "slide4_title": "string",
  "slide4_content": "string",
  "slide5_title": "string",
  "slide5_content": "string",
  "slide6_title": "string",
  "slide6_content": "string",
  "slide7_title": "string",
  "slide7_content": "string"
}}

Content guidelines:
- Be confident, positive, and sales-focused
- Highlight progress and partnership
- Acknowledge risks professionally with mitigation
- End with clear next steps

Key data:
{data_summary}

Risk summary:
High risk items: {risk_data.get("HIGH", 0)}
Medium risk items: {risk_data.get("MEDIUM", 0)}
Low risk items: {risk_data.get("LOW", 0)}

Relevant context:
{rag_context[:6000]}

Generate the JSON now:"""

    try:
        response = await agent.run(task=task)
    except Exception:
        # Silent fallback on error
        return get_fallback_content(customer_name, risk_data)

    # === Extract text from TaskResult ===
    text = ""

    if hasattr(response, "messages") and response.messages:
        for msg in reversed(response.messages):
            if hasattr(msg, "content") and msg.content:
                text = msg.content.strip()
                break
            elif isinstance(msg, dict) and "content" in msg and msg["content"]:
                text = msg["content"].strip()
                break
    elif hasattr(response, "summary") and response.summary:
        text = response.summary.strip()
    elif hasattr(response, "content") and response.content:
        text = response.content.strip()

    if not text:
        return get_fallback_content(customer_name, risk_data)

    # Clean wrappers
    for marker in ["```json", "```", "`", "json\n"]:
        text = text.replace(marker, "")
    text = text.strip()

    # Extract JSON block if needed
    if not text.startswith("{"):
        json_match = re.search(r"\{.*\}", text, re.DOTALL)
        if json_match:
            text = json_match.group(0)
        else:
            return get_fallback_content(customer_name, risk_data)

    try:
        parsed = json.loads(text)

        required_keys = {f"slide{i}_{k}" for i in range(1, 8) for k in ["title", "content"]}
        if required_keys - set(parsed.keys()):
            return get_fallback_content(customer_name, risk_data)

        return parsed

    except json.JSONDecodeError:
        return get_fallback_content(customer_name, risk_data)


def get_fallback_content(customer_name: str, risk_data: dict) -> dict:
    """Clean fallback content"""
    return {
        "slide1_title": f"Strategic Partnership Review - {customer_name}",
        "slide1_content": "Strong, long-term collaboration delivering measurable business value",
        "slide2_title": "Contract Commitments Overview",
        "slide2_content": "All key features clearly defined and actively tracked",
        "slide3_title": "Value Already Delivered",
        "slide3_content": "Multiple critical features live in production\\nSignificant business outcomes achieved",
        "slide4_title": "Active Roadmap & Momentum",
        "slide4_content": "Strong pipeline of planned enhancements\\nClear delivery timeline established",
        "slide5_title": "Risk Management",
        "slide5_content": f"High risks: {risk_data.get('HIGH', 0)}\\nMedium risks: {risk_data.get('MEDIUM', 0)}\\nProactive mitigation plans in place",
        "slide6_title": "Business Impact & ROI",
        "slide6_content": "Realized value through delivered capabilities\\nStrong foundation for continued growth",
        "slide7_title": "Next Steps",
        "slide7_content": "• Schedule executive review meeting\\n• Align on Q2 priorities\\n• Continue close partnership"
    }


def generate_pitch_deck_content_sync(*args, **kwargs):
    """Synchronous wrapper"""
    return asyncio.run(generate_pitch_deck_content(*args, **kwargs))


def build_pptx_from_content(content: dict, customer_name: str) -> str:
    """Build PowerPoint from content"""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = f"{customer_name} Strategic Review"
    slide.placeholders[1].text = f"Prepared {datetime.now().strftime('%B %d, %Y')}"

    # Content slides
    for i in range(1, 8):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.placeholders[1].text_frame
        body.clear()

        title.text = content.get(f"slide{i}_title", f"Slide {i}")

        raw = content.get(f"slide{i}_content", "")
        for line in raw.replace("\\n", "\n").split("\n"):
            if line.strip():
                p = body.add_paragraph()
                p.text = line.strip()
                p.font.size = Pt(18)

    safe_name = "".join(c for c in customer_name if c.isalnum() or c in " _-")
    filename = f"{safe_name}_Pitch_Deck_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    os.makedirs("data", exist_ok=True)
    path = os.path.join("data", filename)
    prs.save(path)
    return path